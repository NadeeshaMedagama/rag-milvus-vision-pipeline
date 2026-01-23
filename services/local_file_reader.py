"""Local file reader service implementation."""
import os
import json
import hashlib
from typing import List, Set, Optional
from pathlib import Path
import docx
import openpyxl
from pptx import Presentation
import pdfplumber

from interfaces import ILocalFileReader, IVisionAnalyzer
from models.data_models import Document, DocumentType


class LocalFileReader(ILocalFileReader):
    """Service for reading local files and directories."""

    # Supported file extensions
    IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.svg', '.webp'}
    DIAGRAM_EXTENSIONS = {'.drawio'}
    DOCUMENT_EXTENSIONS = {'.docx', '.doc'}
    SPREADSHEET_EXTENSIONS = {'.xlsx', '.xls'}
    PDF_EXTENSIONS = {'.pdf'}
    POWERPOINT_EXTENSIONS = {'.pptx', '.ppt'}
    JSON_EXTENSIONS = {'.json'}
    MARKDOWN_EXTENSIONS = {'.md', '.markdown'}
    GRAPHQL_EXTENSIONS = {'.graphql', '.gql'}
    ODT_EXTENSIONS = {'.odt'}
    VIDEO_EXTENSIONS = {'.mp4', '.avi', '.mov', '.mkv', '.webm'}
    EXCALIDRAW_EXTENSIONS = {'.excalidraw'}
    TEXT_EXTENSIONS = {'.txt', '.text', '.log', '.yml', '.yaml', '.xml', '.html', '.htm', '.css', '.js', '.ts', '.py', '.java', '.go', '.rs', '.c', '.cpp', '.h', '.sh', '.bash', '.sql', '.ini', '.cfg', '.conf', '.env'}

    def __init__(self, vision_analyzer: IVisionAnalyzer = None):
        """
        Initialize the local file reader.

        Args:
            vision_analyzer: Optional vision analyzer for processing images
        """
        self.vision_analyzer = vision_analyzer

    @staticmethod
    def normalize_path(file_path: str) -> str:
        """
        Normalize a file path for consistent comparison.

        Args:
            file_path: The file path to normalize

        Returns:
            Normalized file path string
        """
        # Convert to Path object and normalize
        path = Path(file_path)
        # Remove leading ./ if present and normalize
        normalized = str(path).lstrip('./')
        return normalized

    def read_directory(self, directory_path: str, existing_file_paths: Optional[Set[str]] = None) -> List[Document]:
        """
        Read all supported files from a directory recursively.

        Args:
            directory_path: Path to the directory
            existing_file_paths: Optional set of file paths already in vector store to skip

        Returns:
            List of Document objects
        """
        documents = []
        directory = Path(directory_path)

        if not directory.exists():
            print(f"Warning: Directory {directory_path} does not exist")
            return documents

        print(f"Scanning directory: {directory_path}")

        # Normalize existing paths for comparison
        normalized_existing = set()
        if existing_file_paths:
            normalized_existing = {self.normalize_path(p) for p in existing_file_paths}
            print(f"Will skip {len(normalized_existing)} already indexed files")

        skipped_count = 0
        processed_count = 0

        # Walk through all files recursively
        for file_path in directory.rglob('*'):
            if file_path.is_file():
                # Normalize the current file path
                normalized_path = self.normalize_path(str(file_path))

                # Check if file already exists in vector store
                if normalized_path in normalized_existing:
                    skipped_count += 1
                    continue

                try:
                    document = self.read_file(str(file_path))
                    if document:
                        # Store with normalized path for consistent comparison
                        document.file_path = normalized_path
                        documents.append(document)
                        processed_count += 1
                        print(f"Processed: {file_path.name}")
                except Exception as e:
                    print(f"Error processing {file_path}: {str(e)}")

        print(f"\n📊 Local File Processing Summary:")
        print(f"  - Files found: {skipped_count + processed_count}")
        print(f"  - Already indexed (skipped): {skipped_count}")
        print(f"  - Newly processed: {processed_count}")
        return documents

    def read_file(self, file_path: str) -> Document:
        """
        Read a single file.

        Args:
            file_path: Path to the file

        Returns:
            Document object or None if file type not supported
        """
        path = Path(file_path)
        extension = path.suffix.lower()

        # Determine file type and process accordingly
        if extension in self.IMAGE_EXTENSIONS:
            return self._process_image(file_path)
        elif extension in self.DIAGRAM_EXTENSIONS:
            return self._process_diagram(file_path)
        elif extension in self.DOCUMENT_EXTENSIONS:
            return self._process_word_document(file_path)
        elif extension in self.SPREADSHEET_EXTENSIONS:
            return self._process_spreadsheet(file_path)
        elif extension in self.PDF_EXTENSIONS:
            return self._process_pdf(file_path)
        elif extension in self.POWERPOINT_EXTENSIONS:
            return self._process_powerpoint(file_path)
        elif extension in self.JSON_EXTENSIONS:
            return self._process_json(file_path)
        elif extension in self.MARKDOWN_EXTENSIONS:
            return self._process_markdown(file_path)
        elif extension in self.GRAPHQL_EXTENSIONS:
            return self._process_graphql(file_path)
        elif extension in self.ODT_EXTENSIONS:
            return self._process_odt(file_path)
        elif extension in self.VIDEO_EXTENSIONS:
            return self._process_video(file_path)
        elif extension in self.EXCALIDRAW_EXTENSIONS:
            return self._process_excalidraw(file_path)
        elif extension in self.TEXT_EXTENSIONS:
            return self._process_text(file_path)
        else:
            # Skip unsupported file types
            return None

    def _process_image(self, file_path: str) -> Document:
        """Process image files using Vision API."""
        if not self.vision_analyzer:
            return Document(
                content=f"Image file: {os.path.basename(file_path)} (Vision API not configured)",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.IMAGE,
                metadata={"source": "local_directory", "file_type": "image"}
            )

        # Use Vision API to analyze the image
        summary = self.vision_analyzer.generate_summary(file_path)

        return Document(
            content=summary,
            file_path=file_path,
            repository_url="local",
            document_type=DocumentType.IMAGE,
            metadata={
                "source": "local_directory",
                "file_type": "image",
                "analyzed_by": "google_vision_api"
            }
        )

    def _process_diagram(self, file_path: str) -> Document:
        """Process diagram files (.drawio)."""
        # For .drawio files, try to read as XML text
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

            # If there's a corresponding .png file, also analyze it
            png_path = file_path + '.png'
            if os.path.exists(png_path) and self.vision_analyzer:
                vision_summary = self.vision_analyzer.generate_summary(png_path)
                content = f"Diagram File: {os.path.basename(file_path)}\n\n--- Visual Analysis ---\n{vision_summary}\n\n--- Source XML ---\n{content}"
            else:
                content = f"Diagram File: {os.path.basename(file_path)}\n\n{content}"

            return Document(
                content=content,
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.DRAWIO,
                metadata={
                    "source": "local_directory",
                    "file_type": "drawio",
                    "has_png_export": os.path.exists(png_path)
                }
            )
        except Exception as e:
            return Document(
                content=f"Diagram file: {os.path.basename(file_path)} (Error reading: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.DRAWIO,
                metadata={"source": "local_directory", "file_type": "drawio", "error": str(e)}
            )

    def _process_word_document(self, file_path: str) -> Document:
        """Process Word documents (.docx, .doc)."""
        try:
            # Only .docx is supported by python-docx
            if file_path.endswith('.docx'):
                doc = docx.Document(file_path)
                content = []

                # Extract paragraphs
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        content.append(paragraph.text)

                # Extract tables
                for table in doc.tables:
                    for row in table.rows:
                        row_text = ' | '.join([cell.text for cell in row.cells])
                        if row_text.strip():
                            content.append(row_text)

                full_content = f"Word Document: {os.path.basename(file_path)}\n\n" + "\n".join(content)

                return Document(
                    content=full_content,
                    file_path=file_path,
                    repository_url="local",
                    document_type=DocumentType.WORD_DOCUMENT,
                    metadata={
                        "source": "local_directory",
                        "file_type": "word_document",
                        "paragraph_count": len(doc.paragraphs),
                        "table_count": len(doc.tables)
                    }
                )
            else:
                return Document(
                    content=f"Word document: {os.path.basename(file_path)} (.doc format not supported, only .docx)",
                    file_path=file_path,
                    repository_url="local",
                    document_type=DocumentType.WORD_DOCUMENT,
                    metadata={"source": "local_directory", "file_type": "word_document"}
                )
        except Exception as e:
            return Document(
                content=f"Word document: {os.path.basename(file_path)} (Error reading: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.WORD_DOCUMENT,
                metadata={"source": "local_directory", "file_type": "word_document", "error": str(e)}
            )

    def _process_spreadsheet(self, file_path: str) -> Document:
        """Process spreadsheet files (.xlsx, .xls)."""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            content = [f"Spreadsheet: {os.path.basename(file_path)}\n"]

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content.append(f"\n--- Sheet: {sheet_name} ---")

                # Read up to 100 rows to avoid too much data
                max_rows = min(sheet.max_row, 100)
                for row_idx, row in enumerate(sheet.iter_rows(max_row=max_rows, values_only=True), 1):
                    row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip():
                        content.append(row_text)

                if sheet.max_row > 100:
                    content.append(f"... (Truncated, total rows: {sheet.max_row})")

            full_content = "\n".join(content)

            return Document(
                content=full_content,
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.SPREADSHEET,
                metadata={
                    "source": "local_directory",
                    "file_type": "spreadsheet",
                    "sheet_count": len(workbook.sheetnames),
                    "sheet_names": workbook.sheetnames
                }
            )
        except Exception as e:
            return Document(
                content=f"Spreadsheet: {os.path.basename(file_path)} (Error reading: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.SPREADSHEET,
                metadata={"source": "local_directory", "file_type": "spreadsheet", "error": str(e)}
            )

    def _process_pdf(self, file_path: str) -> Document:
        """Process PDF files using pdfplumber."""
        try:
            content = [f"PDF Document: {os.path.basename(file_path)}\n"]
            page_count = 0

            with pdfplumber.open(file_path) as pdf:
                page_count = len(pdf.pages)
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        content.append(f"\n--- Page {i + 1} ---\n{page_text}")

                    # Extract tables if present
                    tables = page.extract_tables()
                    for table_idx, table in enumerate(tables):
                        if table:
                            content.append(f"\n[Table {table_idx + 1} on Page {i + 1}]")
                            for row in table:
                                row_text = ' | '.join([str(cell) if cell else '' for cell in row])
                                if row_text.strip():
                                    content.append(row_text)

            full_content = "\n".join(content)

            # If no text was extracted, try to use Vision API for scanned PDFs
            if len(full_content.strip()) < 50 and self.vision_analyzer:
                # PDF might be image-based/scanned
                full_content += "\n\n[Note: PDF appears to be image-based. Text extraction limited.]"

            return Document(
                content=full_content,
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.PDF,
                metadata={
                    "source": "local_directory",
                    "file_type": "pdf",
                    "page_count": page_count
                }
            )
        except Exception as e:
            return Document(
                content=f"PDF Document: {os.path.basename(file_path)} (Error reading: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.PDF,
                metadata={"source": "local_directory", "file_type": "pdf", "error": str(e)}
            )

    def _process_powerpoint(self, file_path: str) -> Document:
        """Process PowerPoint files (.pptx, .ppt)."""
        try:
            # Only .pptx is supported by python-pptx
            if file_path.endswith('.pptx'):
                prs = Presentation(file_path)
                content = [f"PowerPoint Presentation: {os.path.basename(file_path)}\n"]
                slide_count = len(prs.slides)

                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_content = []

                    for shape in slide.shapes:
                        # Extract text from shapes
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_content.append(shape.text)

                        # Extract text from tables
                        if shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                row_text = ' | '.join([cell.text for cell in row.cells])
                                if row_text.strip():
                                    slide_content.append(row_text)

                    if slide_content:
                        content.append(f"\n--- Slide {slide_num} ---")
                        content.extend(slide_content)

                full_content = "\n".join(content)

                return Document(
                    content=full_content,
                    file_path=file_path,
                    repository_url="local",
                    document_type=DocumentType.POWERPOINT,
                    metadata={
                        "source": "local_directory",
                        "file_type": "powerpoint",
                        "slide_count": slide_count
                    }
                )
            else:
                # .ppt format is not supported
                return Document(
                    content=f"PowerPoint: {os.path.basename(file_path)} (.ppt format not supported, only .pptx)",
                    file_path=file_path,
                    repository_url="local",
                    document_type=DocumentType.POWERPOINT,
                    metadata={"source": "local_directory", "file_type": "powerpoint", "unsupported_format": True}
                )
        except Exception as e:
            return Document(
                content=f"PowerPoint: {os.path.basename(file_path)} (Error reading: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.POWERPOINT,
                metadata={"source": "local_directory", "file_type": "powerpoint", "error": str(e)}
            )

    def _process_json(self, file_path: str) -> Document:
        """Process JSON files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                data = json.load(f)

            # Pretty print the JSON for better readability
            formatted_json = json.dumps(data, indent=2, ensure_ascii=False)
            content = f"JSON File: {os.path.basename(file_path)}\n\n{formatted_json}"

            return Document(
                content=content,
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.JSON,
                metadata={
                    "source": "local_directory",
                    "file_type": "json",
                    "keys": list(data.keys()) if isinstance(data, dict) else None,
                    "is_array": isinstance(data, list),
                    "item_count": len(data) if isinstance(data, (list, dict)) else None
                }
            )
        except json.JSONDecodeError as e:
            # If JSON is invalid, read as plain text
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return Document(
                content=f"JSON File (invalid JSON): {os.path.basename(file_path)}\n\n{content}",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.JSON,
                metadata={"source": "local_directory", "file_type": "json", "error": f"Invalid JSON: {str(e)}"}
            )
        except Exception as e:
            return Document(
                content=f"JSON File: {os.path.basename(file_path)} (Error reading: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.JSON,
                metadata={"source": "local_directory", "file_type": "json", "error": str(e)}
            )

    def _process_markdown(self, file_path: str) -> Document:
        """Process Markdown files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

            full_content = f"Markdown Document: {os.path.basename(file_path)}\n\n{content}"

            return Document(
                content=full_content,
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.MARKDOWN,
                metadata={
                    "source": "local_directory",
                    "file_type": "markdown",
                    "char_count": len(content),
                    "line_count": content.count('\n') + 1
                }
            )
        except Exception as e:
            return Document(
                content=f"Markdown Document: {os.path.basename(file_path)} (Error reading: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.MARKDOWN,
                metadata={"source": "local_directory", "file_type": "markdown", "error": str(e)}
            )

    def _process_graphql(self, file_path: str) -> Document:
        """Process GraphQL schema files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

            full_content = f"GraphQL Schema: {os.path.basename(file_path)}\n\n{content}"

            return Document(
                content=full_content,
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.GRAPHQL,
                metadata={
                    "source": "local_directory",
                    "file_type": "graphql",
                    "char_count": len(content),
                    "has_query": "type Query" in content or "query" in content.lower(),
                    "has_mutation": "type Mutation" in content or "mutation" in content.lower()
                }
            )
        except Exception as e:
            return Document(
                content=f"GraphQL Schema: {os.path.basename(file_path)} (Error reading: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.GRAPHQL,
                metadata={"source": "local_directory", "file_type": "graphql", "error": str(e)}
            )

    def _process_odt(self, file_path: str) -> Document:
        """Process OpenDocument Text files (.odt)."""
        try:
            from odf import text as odf_text
            from odf.opendocument import load as odf_load

            doc = odf_load(file_path)
            content = []

            # Extract all text from paragraphs
            for paragraph in doc.getElementsByType(odf_text.P):
                para_text = ""
                for node in paragraph.childNodes:
                    if node.nodeType == node.TEXT_NODE:
                        para_text += str(node)
                    elif hasattr(node, 'childNodes'):
                        for child in node.childNodes:
                            if child.nodeType == child.TEXT_NODE:
                                para_text += str(child)
                if para_text.strip():
                    content.append(para_text)

            full_content = f"OpenDocument Text: {os.path.basename(file_path)}\n\n" + "\n".join(content)

            return Document(
                content=full_content,
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.ODT,
                metadata={
                    "source": "local_directory",
                    "file_type": "odt",
                    "paragraph_count": len(content)
                }
            )
        except ImportError:
            return Document(
                content=f"OpenDocument Text: {os.path.basename(file_path)} (odfpy library not installed. Run: pip install odfpy)",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.ODT,
                metadata={"source": "local_directory", "file_type": "odt", "error": "odfpy not installed"}
            )
        except Exception as e:
            return Document(
                content=f"OpenDocument Text: {os.path.basename(file_path)} (Error reading: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.ODT,
                metadata={"source": "local_directory", "file_type": "odt", "error": str(e)}
            )

    def _process_video(self, file_path: str) -> Document:
        """Process video files (metadata only, no content extraction)."""
        try:
            file_size = os.path.getsize(file_path)
            file_size_mb = round(file_size / (1024 * 1024), 2)

            content = f"Video File: {os.path.basename(file_path)}\n\n"
            content += f"File Size: {file_size_mb} MB\n"
            content += f"File Path: {file_path}\n"
            content += "\n[Note: Video content cannot be directly extracted. This is metadata only.]"

            return Document(
                content=content,
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.VIDEO,
                metadata={
                    "source": "local_directory",
                    "file_type": "video",
                    "file_size_bytes": file_size,
                    "file_size_mb": file_size_mb
                }
            )
        except Exception as e:
            return Document(
                content=f"Video File: {os.path.basename(file_path)} (Error reading metadata: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.VIDEO,
                metadata={"source": "local_directory", "file_type": "video", "error": str(e)}
            )

    def _process_excalidraw(self, file_path: str) -> Document:
        """Process Excalidraw diagram files (JSON-based)."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                data = json.load(f)

            content = [f"Excalidraw Diagram: {os.path.basename(file_path)}\n"]

            # Extract text elements from the diagram
            elements = data.get('elements', [])
            text_elements = []
            shape_counts = {}

            for element in elements:
                elem_type = element.get('type', 'unknown')
                shape_counts[elem_type] = shape_counts.get(elem_type, 0) + 1

                # Extract text content
                if elem_type == 'text':
                    text = element.get('text', '').strip()
                    if text:
                        text_elements.append(text)

            if text_elements:
                content.append("\n--- Text Content ---")
                for text in text_elements:
                    content.append(f"• {text}")

            content.append(f"\n--- Diagram Statistics ---")
            content.append(f"Total Elements: {len(elements)}")
            for shape_type, count in sorted(shape_counts.items()):
                content.append(f"  • {shape_type}: {count}")

            full_content = "\n".join(content)

            return Document(
                content=full_content,
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.EXCALIDRAW,
                metadata={
                    "source": "local_directory",
                    "file_type": "excalidraw",
                    "element_count": len(elements),
                    "text_element_count": len(text_elements),
                    "shape_types": shape_counts
                }
            )
        except json.JSONDecodeError as e:
            return Document(
                content=f"Excalidraw Diagram: {os.path.basename(file_path)} (Invalid JSON: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.EXCALIDRAW,
                metadata={"source": "local_directory", "file_type": "excalidraw", "error": f"Invalid JSON: {str(e)}"}
            )
        except Exception as e:
            return Document(
                content=f"Excalidraw Diagram: {os.path.basename(file_path)} (Error reading: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.EXCALIDRAW,
                metadata={"source": "local_directory", "file_type": "excalidraw", "error": str(e)}
            )

    def _process_text(self, file_path: str) -> Document:
        """Process plain text and code files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

            extension = Path(file_path).suffix.lower()
            file_type_map = {
                '.txt': 'Plain Text',
                '.text': 'Plain Text',
                '.log': 'Log File',
                '.yml': 'YAML',
                '.yaml': 'YAML',
                '.xml': 'XML',
                '.html': 'HTML',
                '.htm': 'HTML',
                '.css': 'CSS',
                '.js': 'JavaScript',
                '.ts': 'TypeScript',
                '.py': 'Python',
                '.java': 'Java',
                '.go': 'Go',
                '.rs': 'Rust',
                '.c': 'C',
                '.cpp': 'C++',
                '.h': 'C/C++ Header',
                '.sh': 'Shell Script',
                '.bash': 'Bash Script',
                '.sql': 'SQL',
                '.ini': 'INI Config',
                '.cfg': 'Config File',
                '.conf': 'Config File',
                '.env': 'Environment File'
            }

            file_type = file_type_map.get(extension, 'Text File')
            full_content = f"{file_type}: {os.path.basename(file_path)}\n\n{content}"

            return Document(
                content=full_content,
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.TEXT,
                metadata={
                    "source": "local_directory",
                    "file_type": file_type.lower().replace(' ', '_'),
                    "extension": extension,
                    "char_count": len(content),
                    "line_count": content.count('\n') + 1
                }
            )
        except Exception as e:
            return Document(
                content=f"Text File: {os.path.basename(file_path)} (Error reading: {str(e)})",
                file_path=file_path,
                repository_url="local",
                document_type=DocumentType.TEXT,
                metadata={"source": "local_directory", "file_type": "text", "error": str(e)}
            )

